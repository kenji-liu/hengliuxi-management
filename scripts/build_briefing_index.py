"""
把金質獎評審簡報（.pptx）逐頁轉成可檢索的知識庫
======================================================================
評審委員問答準備手冊大量以頁碼引用簡報（P.16、P.61、P.121…），因此
AI 答詢必須能查到「該頁實際寫了什麼」，而不是只有手冊的二手摘要。

以「一頁投影片＝一個知識單元」儲存，頁碼即投影片編號，與手冊引用一致。
簡報檔本身有數百 MB（內嵌大量照片），但純文字僅數萬字，適合直接全文檢索。

產出：webapp/data/briefing_slides.json

用法：
    python scripts/build_briefing_index.py [簡報.pptx]
"""

from __future__ import annotations

import glob
import json
import os
import re
import sys

OUT_PATH = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                        "webapp", "data", "briefing_slides.json")

# 依簡報章節標題推斷所屬評分構面，供答詢時對應委員關注面向
SECTION_PATTERNS = [
    (r"維護管理制度", "維護管理制度"),
    (r"維護作業品質|維護作業", "維護作業品質"),
    (r"文件管理|技術轉移", "維護文件管理"),
    (r"節能減碳|碳排|碳匯", "節能減碳"),
    (r"防災|安全|勞安", "防災與安全"),
    (r"環境保育|生態|棲地|植生", "環境保育"),
    (r"創新科技|AI|智慧|數位", "創新科技"),
]


def shape_texts(shapes) -> list:
    """遞迴抽出所有文字，含群組與表格。"""
    out = []
    for shape in shapes:
        try:
            # 群組圖形（shape_type 6）需遞迴進入
            if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):
                out += shape_texts(shape.shapes)
                continue
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    out.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                    if any(cells):
                        out.append(" | ".join(cells))
        except Exception:
            continue
    return out


def guess_section(text: str, current: str) -> str:
    for pattern, label in SECTION_PATTERNS:
        if re.search(pattern, text):
            return label
    return current


def build(src: str) -> dict:
    from pptx import Presentation
    deck = Presentation(src)

    slides = []
    section = ""
    for index, slide in enumerate(deck.slides, 1):
        blocks = shape_texts(slide.shapes)
        body = "\n".join(blocks).strip()

        notes = ""
        try:
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            notes = ""

        # 章節頁通常只有一行標題，用來標記後續頁面所屬構面
        if body and len(blocks) <= 2 and len(body) <= 30:
            section = guess_section(body, section)
        else:
            section = guess_section(body[:120], section)

        if not body and not notes:
            continue

        slides.append({
            "page": index,
            "title": blocks[0][:60] if blocks else "",
            "section": section,
            "text": body[:2500],
            "notes": notes[:800],
        })

    return {
        "source": os.path.basename(src),
        "totalSlides": len(deck.slides),
        "indexedSlides": len(slides),
        "slides": slides,
    }


def main() -> int:
    if len(sys.argv) > 1:
        src = sys.argv[1]
    else:
        pattern = os.path.join(os.path.expanduser("~"), "Downloads", "*金質獎評審簡報*.pptx")
        found = sorted(glob.glob(pattern))
        if not found:
            print("請指定簡報檔路徑")
            return 1
        src = found[-1]

    if not os.path.exists(src):
        print(f"找不到檔案：{src}")
        return 1

    print(f"讀取：{os.path.basename(src)}")
    data = build(src)
    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    with open(OUT_PATH, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=1)

    chars = sum(len(s["text"]) for s in data["slides"])
    print(f"已產出：{OUT_PATH}")
    print(f"  投影片 {data['totalSlides']} 頁，收錄 {data['indexedSlides']} 頁，共 {chars:,} 字")
    sections = {}
    for s in data["slides"]:
        key = s["section"] or "（未分類）"
        sections[key] = sections.get(key, 0) + 1
    for key, count in sorted(sections.items(), key=lambda x: -x[1]):
        print(f"  {count:>3} 頁  {key}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
